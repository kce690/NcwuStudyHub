from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from image_exporter import export_picture_shape


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in _iter_shapes(shape.shapes):
                yield sub_shape


def _extract_paragraphs(shape) -> list[dict]:
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return []

    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = " ".join((paragraph.text or "").split())
        if not text:
            continue
        paragraphs.append(
            {
                "text": text,
                "level": int(getattr(paragraph, "level", 0) or 0),
            }
        )
    return paragraphs


def extract_pptx_content(pptx_path: Path, images_dir: Path, logger) -> tuple[list[dict], dict]:
    """
    提取 PPTX 内容：
    - 每页标题、文本块、项目符号
    - 每页图片导出并记录相对路径
    """
    presentation = Presentation(str(pptx_path))
    slides_data: list[dict] = []
    total_images = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = ""
        try:
            title_shape = slide.shapes.title
            if title_shape and title_shape.text:
                title = " ".join(title_shape.text.split())
        except Exception:  # noqa: BLE001
            title = ""

        text_blocks: list[str] = []
        bullet_points: list[dict] = []
        image_paths: list[str] = []

        image_index = 1
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                rel_image = export_picture_shape(shape, slide_number, image_index, images_dir, logger)
                if rel_image:
                    image_paths.append(rel_image)
                    total_images += 1
                    image_index += 1

            paragraphs = _extract_paragraphs(shape)
            if not paragraphs:
                continue

            block_lines = [p["text"] for p in paragraphs]
            block_text = "\n".join(block_lines).strip()
            if block_text and block_text != title:
                text_blocks.append(block_text)

            for p in paragraphs:
                bullet_points.append({"level": p["level"], "text": p["text"]})

        if not title and text_blocks:
            title = text_blocks[0].splitlines()[0][:80]

        slide_item = {
            "slide_number": slide_number,
            "title": title or f"第{slide_number}页",
            "text_blocks": text_blocks,
            "bullet_points": bullet_points,
            "image_paths": image_paths,
        }
        if not text_blocks and image_paths:
            slide_item["note"] = "该页主要为图片内容"

        slides_data.append(slide_item)

    stats = {
        "slide_count": len(slides_data),
        "image_count": total_images,
    }
    logger.info("提取完成：slide=%s image=%s", stats["slide_count"], stats["image_count"])
    return slides_data, stats
